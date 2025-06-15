import os
import re
import itertools
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def add_slide_with_image(ppt, image_path, title, description):
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(1.5))
    textbox.text_frame.add_paragraph().text = description

def sanitize_filename(filename):
    filename = re.sub(r'[\\/*?:"<>|]', "", filename)
    return filename.replace(" ", "_")

def load_or_create_ppt(ppt_path):
    if os.path.exists(ppt_path):
        return Presentation(ppt_path)
    return Presentation()

def save_plot_with_directory_check(image_path, plt_obj):
    directory = os.path.dirname(image_path)
    if not os.path.exists(directory):
        os.makedirs(directory)
    plt_obj.savefig(image_path)
    plt_obj.close()

def perform_descriptive_statistics(data, query, ppt_path, sanitized_query):
    ppt = load_or_create_ppt(ppt_path)  # Load or create PowerPoint presentation
    chartnum = 0
    print("\nDescriptive Statistics:")

    # Compute general statistics for the whole dataset
    desc_stats = data.describe(include='all').transpose()
    print(desc_stats)

    # Identify all categorical columns
    categorical_cols = data.select_dtypes(include=['object']).columns

    if len(categorical_cols) < 2:
        print("Not enough categorical columns for combination analysis.")
    else:
        for fixed_col in categorical_cols:
            varying_cols = [col for col in categorical_cols if col != fixed_col]
            for var_pair in itertools.combinations(varying_cols, 2):
                col1, col2 = var_pair
                grouped_data = data.groupby([fixed_col, col1, col2]).size().reset_index(name='count')

                for fixed_value in grouped_data[fixed_col].unique():
                    subset = grouped_data[grouped_data[fixed_col] == fixed_value]

                    plt.figure(figsize=(12, 8))
                    sns.barplot(x=subset[col1] + " - " + subset[col2], y=subset['count'], palette='Set2')

                    plt.title(f'Variation of {col1} & {col2} when {fixed_col} is {fixed_value}')
                    plt.xlabel(f'{col1} & {col2} Combinations')
                    plt.ylabel('Count')
                    plt.xticks(rotation=45, ha='right')
                    plt.tight_layout()

                    chartnum += 1
                    image_filename = f"variation_{chartnum}.png"
                    image_path = os.path.join("output", image_filename)
                    save_plot_with_directory_check(image_path, plt)

                    description = f"This chart shows how {col1} and {col2} vary when {fixed_col} is {fixed_value}."
                    add_slide_with_image(ppt, image_path, f"{col1} & {col2} variation by {fixed_col}", description)

    numerical_cols = data.select_dtypes(include=['number']).columns
    for col in numerical_cols:
        print(f"\n--- Descriptive statistics for numerical column '{col}' ---")
        stats = data[col].describe()
        print(stats)

        plt.figure(figsize=(10, 6))
        sns.histplot(data[col], kde=True, bins=30)
        plt.title(f"Distribution of '{col}'")
        plt.xlabel(col)
        plt.ylabel("Frequency")
        image_path = os.path.join("output", f"{col}_distribution.png")
        
        save_plot_with_directory_check(image_path, plt)
        description = f"This chart shows the distribution of values in the '{col}' numerical column."
        add_slide_with_image(ppt, image_path, f"Distribution of '{col}'", description)

    ppt.save(ppt_path)
    print("PowerPoint presentation created/updated successfully.")
    return ppt_path
