### OLD VERSION NO COMPONENT SEPERATION ###


import os
import re
import itertools
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
import torch
import nltk
from nltk.corpus import stopwords
from flask import Flask, request, jsonify, send_from_directory
from transformers import BertTokenizer, BertModel, pipeline
from pptx import Presentation
from pptx.util import Inches
import traceback
from flask_cors import CORS


# Download NLTK stopwords
nltk.download('stopwords')

# Create Flask app and configure folders
app = Flask(__name__)
CORS(app)
app.config["UPLOAD_FOLDER"] = os.path.join(os.getcwd(), "uploads")
app.config["OUTPUT_FOLDER"] = os.path.join(os.getcwd(), "outputs")

# Create folders if they do not exist
for folder in [app.config["UPLOAD_FOLDER"], app.config["OUTPUT_FOLDER"]]:
    if not os.path.exists(folder):
        os.makedirs(folder)

# Load pre-trained BERT components (globally)
tokenizer_bert = BertTokenizer.from_pretrained("bert-base-uncased")
model_bert = BertModel.from_pretrained("bert-base-uncased")

def preprocess_data(data):
    data = data.replace(r'^\s*$', np.nan, regex=True)
    data = data.dropna(axis=1, how='all')
    for col in data.columns:
        if data[col].dtype in ['float64', 'int64']:
            data[col] = data[col].fillna(data[col].median())
        elif data[col].dtype == 'object':
            data[col] = data[col].fillna(data[col].mode()[0])
    return data

def get_bert_embeddings(texts):
    inputs = tokenizer_bert(texts, return_tensors="pt", padding=True, truncation=True, max_length=512)
    with torch.no_grad():
        outputs = model_bert(**inputs)
    embeddings = outputs.last_hidden_state[:, 0, :]
    return embeddings

def generate_tags_with_context(data):
    tags = {}
    for col in data.columns:
        column_values = data[col].dropna().astype(str).unique()
        generalized_tags = {col.lower(), col.replace("_", " ").lower()}
        for value in column_values:
            value_tags = {value.lower(), value.replace("_", " ").lower()}
            generalized_tags.update(value_tags)
        generalized_tags = list(generalized_tags)
        embeddings = get_bert_embeddings(generalized_tags)
        tags[col] = {
            "tags": generalized_tags,
            "embeddings": embeddings,
            "candidates": column_values
        }
    return tags

def extract_dynamic_tags(query):
    stop_words = set(stopwords.words('english'))
    words = set(re.findall(r'\b\w+\b', query.lower()))
    meaningful_words = words - stop_words
    return meaningful_words

def match_query_to_tags(query, tag_embeddings, threshold=0.5):
    query_embedding = get_bert_embeddings([query])[0]
    matched_columns = []
    matched_values = {}
    additional_tags = extract_dynamic_tags(query)
    all_tags = set()
    for col, tag_data in tag_embeddings.items():
        column_tags = set(tag_data["tags"])
        if column_tags & additional_tags:
            all_tags.add(col)
    for col, tag_data in tag_embeddings.items():
        similarities = torch.cosine_similarity(query_embedding, tag_data["embeddings"], dim=1)
        best_score = torch.max(similarities).item()
        if best_score > threshold and col in all_tags:
            matched_columns.append(col)
            for idx, value in enumerate(tag_data["candidates"]):
                value_similarity = similarities[idx].item()
                if value_similarity > threshold:
                    if col not in matched_values:
                        matched_values[col] = []
                    matched_values[col].append((value, value_similarity))
    return matched_columns, matched_values

def identify_intent(query):
    intents = ["retrieve", "comparison", "analysis", "filter", "summary"]
    zero_shot_classifier = pipeline("zero-shot-classification", model="facebook/bart-large-mnli")
    result = zero_shot_classifier(query, candidate_labels=intents)
    return result['labels'][0]

def extract_entities(query):
    entities = {}
    ner_pipeline = pipeline("ner", model="dslim/bert-large-NER", tokenizer="dslim/bert-large-NER")
    ner_results = ner_pipeline(query)
    for result in ner_results:
        entity_type = result['entity']
        entity_value = result['word']
        if entity_type not in entities:
            entities[entity_type] = []
        entities[entity_type].append(entity_value)
    return entities

def extract_conditions_from_query(query):
    conditions = {}
    between_matches = re.findall(r"(\w+)\s+between\s+(\d+)\s+and\s+(\d+)", query, re.IGNORECASE)
    for match in between_matches:
        column, lower, upper = match
        conditions[column.lower()] = {"operator": "between", "value": [int(lower), int(upper)]}
    other_matches = re.findall(r"(\w+)\s+(above|below|greater than|less than|>=|<=|>|<|=)\s+(\d+)", query, re.IGNORECASE)
    for match in other_matches:
        column, operator, value = match
        operator = operator.lower().replace("above", ">").replace("below", "<").replace("greater than", ">").replace("less than", "<")
        conditions[column.lower()] = {"operator": operator, "value": int(value)}
    return conditions

def apply_conditions(data, conditions):
    filtered_data = data.copy()
    for column, condition in conditions.items():
        if column in data.columns:
            operator = condition["operator"]
            value = condition["value"]
            if operator == ">":
                filtered_data = filtered_data[filtered_data[column] > value]
            elif operator == "<":
                filtered_data = filtered_data[filtered_data[column] < value]
            elif operator == ">=":
                filtered_data = filtered_data[filtered_data[column] >= value]
            elif operator == "<=":
                filtered_data = filtered_data[filtered_data[column] <= value]
            elif operator == "=":
                filtered_data = filtered_data[filtered_data[column] == value]
            elif operator == "between":
                lower, upper = value
                filtered_data = filtered_data[(filtered_data[column] >= lower) & (filtered_data[column] <= upper)]
    return filtered_data

def filter_rows_based_on_common_keywords(data, matched_values, query):
    filtered_data = data.copy()
    query_keywords = extract_dynamic_tags(query)
    combined_keywords = set(query_keywords)
    for col, matches in matched_values.items():
        combined_keywords.update([word for word, score in matches])
    common_keywords = query_keywords.intersection(combined_keywords)
    if not common_keywords:
        print("No common keywords found between the query and dataset values.")
        return data
    print(f"Common Keywords Found: {common_keywords}")
    for col in matched_values.keys():
        if col in filtered_data.columns:
            pattern = '|'.join(re.escape(word) for word in common_keywords)
            filtered_data = filtered_data[filtered_data[col].astype(str).str.contains(pattern, case=False, na=False)]
    return filtered_data

def add_slide_with_image(ppt, image_path, title, description):
    slide_layout = ppt.slide_layouts[5]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(9))
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(1.5))
    textbox.text_frame.add_paragraph().text = description

def sanitize_filename(filename):
    filename = re.sub(r'[\\/*?:"<>|]', "", filename)
    return filename.replace(" ", "_")

def load_or_create_ppt(ppt_path):
    if os.path.exists(ppt_path):
        return Presentation(ppt_path)
    return Presentation()

def save_plot_with_directory_check(image_path, plt_obj):
    directory = os.path.dirname(image_path)
    if not os.path.exists(directory):
        os.makedirs(directory)
    plt_obj.savefig(image_path)
    plt_obj.close()

# def perform_descriptive_statistics(data, query, ppt_path, sanitized_query):
#     ppt = load_or_create_ppt(ppt_path)
#     categorical_cols = data.select_dtypes(include=['object']).columns
#     column_combinations = []
#     for r in range(2, len(categorical_cols) + 1):
#         column_combinations += list(itertools.combinations(categorical_cols, r))
#     for columns in column_combinations:
#         group_column = data[list(columns)].astype(str).agg(" - ".join, axis=1)
#         grouped_data = group_column.value_counts().reset_index(name='count')
#         grouped_data.columns = ['Combination', 'count']
#         grouped_data = grouped_data.sort_values(by='count', ascending=False)
#         plt.figure(figsize=(12, 8))
#         sns.barplot(x='Combination', y='count', data=grouped_data, palette='Set2')
#         plt.title(f'Aims at answering {query}')
#         plt.xlabel(f'Combinations of {" & ".join(columns)}')
#         plt.ylabel('Count')
#         plt.xticks(rotation=45, ha='right')
#         plt.tight_layout()
#         image_path = os.path.join(app.config["OUTPUT_FOLDER"], f"grouped_bar_chart_{sanitized_query}.png")
#         save_plot_with_directory_check(image_path, plt)
#         add_slide_with_image(ppt, image_path, f"{query}", f"Grouped bar chart for {columns}")
#     numerical_cols = data.select_dtypes(include=['number']).columns
#     for col in numerical_cols:
#         plt.figure(figsize=(10, 6))
#         sns.histplot(data[col], kde=True, bins=30)
#         plt.title(f"Distribution of '{col}'")
#         plt.xlabel(col)
#         plt.ylabel("Frequency")
#         image_path = os.path.join(app.config["OUTPUT_FOLDER"], f"{col}_distribution.png")
#         save_plot_with_directory_check(image_path, plt)
#         add_slide_with_image(ppt, image_path, f"Distribution of '{col}'", f"Histogram for {col}")
#     ppt.save(ppt_path)
#     return ppt_path

def perform_descriptive_statistics(data, query, ppt_path, sanitized_query):

    ppt = load_or_create_ppt(ppt_path)  # Load or create PowerPoint presentation
    chartnum = 0
    print("\nDescriptive Statistics:")

    # Compute general statistics for the whole dataset
    desc_stats = data.describe(include='all').transpose()
    print(desc_stats)

    # Identify all categorical columns
    categorical_cols = data.select_dtypes(include=['object']).columns

    # If fewer than 3 categorical columns, skip advanced analysis
    if len(categorical_cols) < 3:
        print("Not enough categorical columns for combination analysis.")
    else:
        # Loop through each column and fix it while analyzing the variation in other two
        for fixed_col in categorical_cols:
            varying_cols = [col for col in categorical_cols if col != fixed_col]

            # Generate combinations of remaining columns (pairs)
            for var_pair in itertools.combinations(varying_cols, 2):
                col1, col2 = var_pair

                # Group by the selected columns
                grouped_data = data.groupby([fixed_col, col1, col2]).size().reset_index(name='count')

                for fixed_value in grouped_data[fixed_col].unique():
                    subset = grouped_data[grouped_data[fixed_col] == fixed_value]

                    plt.figure(figsize=(12, 8))
                    sns.barplot(x=subset[col1] + " - " + subset[col2], y=subset['count'], palette='Set2')

                    plt.title(f'Variation of {col1} & {col2} when {fixed_col} is {fixed_value}')
                    plt.xlabel(f'{col1} & {col2} Combinations')
                    plt.ylabel('Count')
                    plt.xticks(rotation=45, ha='right')
                    plt.tight_layout()

                    # Save the plot
                    chartnum += 1
                    image_filename = f"variation_{chartnum}.png"
                    image_path = os.path.join(app.config["OUTPUT_FOLDER"], image_filename)
                    save_plot_with_directory_check(image_path, plt)

                    # Add the slide to PowerPoint
                    description = f"This chart shows how {col1} and {col2} vary when {fixed_col} is {fixed_value}."
                    add_slide_with_image(ppt, image_path, f"{col1} & {col2} variation by {fixed_col}", description)

    # Handling numerical columns
    numerical_cols = data.select_dtypes(include=['number']).columns
    for col in numerical_cols:
        print(f"\n--- Descriptive statistics for numerical column '{col}' ---")
        
        # Compute the statistics for numerical columns
        stats = data[col].describe()
        print(stats)
        
        # Plot the distribution of the numerical column
        plt.figure(figsize=(10, 6))
        sns.histplot(data[col], kde=True, bins=30)
        plt.title(f"Distribution of '{col}'")
        plt.xlabel(col)
        plt.ylabel("Frequency")
        image_path = os.path.join(app.config["OUTPUT_FOLDER"], f"{col}_distribution.png")
        
        save_plot_with_directory_check(image_path, plt)

        # Add distribution plot and description to the PPT
        description = f"This chart shows the distribution of values in the '{col}' numerical column."
        add_slide_with_image(ppt, image_path, f"Distribution of '{col}'", description)

    # Save the PowerPoint file (overwrites if it already exists)
    ppt.save(ppt_path)
    print("PowerPoint presentation created/updated successfully.")
    return ppt_path


# ----------------- End of Functions -----------------
@app.route("/", methods=["GET"])
def index():
    return "Flask API is running. Please use the front end to interact with it."

@app.route("/", methods=["POST"])
def process_file():
    try:
        if "file" not in request.files or request.form.get("query", "") == "":
            return jsonify({"error": "Please upload an Excel file and enter a query."}), 400
        
        file = request.files["file"]
        query = request.form.get("query")
        if file.filename == "":
            return jsonify({"error": "No selected file"}), 400
        
        filename = file.filename
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], filename)
        file.save(file_path)
        
        data = pd.read_excel(file_path)
        cleaned_data = preprocess_data(data)
        conditions = extract_conditions_from_query(query)
        entities = extract_entities(query)
        user_intent = identify_intent(query)
        tag_embeddings = generate_tags_with_context(cleaned_data)
        matched_columns, matched_values = match_query_to_tags(query, tag_embeddings, threshold=0.7)
        
        if matched_columns:
            filtered_data = apply_conditions(cleaned_data[matched_columns], conditions)
        else:
            filtered_data = cleaned_data
        final_filtered_data = filter_rows_based_on_common_keywords(filtered_data, matched_values, query)
        
        sanitized_query = sanitize_filename(query)
        ppt_filename = f"{sanitized_query}.pptx"
        ppt_path = os.path.join(app.config["OUTPUT_FOLDER"], ppt_filename)
        perform_descriptive_statistics(filtered_data, query, ppt_path, sanitized_query)
        
        download_url = f"http://127.0.0.1:5000/download/{ppt_filename}"
        return jsonify({"downloadUrl": download_url})
    
    except Exception as e:
    
        error_message = str(e)
        traceback_details = traceback.format_exc()
        print(f"Error: {error_message}\n{traceback_details}")  # error logging
        return jsonify({"error": f"Error: {error_message}"}), 500


@app.route("/download/<filename>")
def download_file(filename):
    return send_from_directory(app.config["OUTPUT_FOLDER"], filename, as_attachment=True)

if __name__ == "__main__":
    import webbrowser
    port = 5000
    # url = f"http://127.0.0.1:{port}/"
    # webbrowser.open(url)
    app.run(debug=True, port=port)
